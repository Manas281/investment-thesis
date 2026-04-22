from flask import Flask, request, jsonify, send_file
from pptx import Presentation
import os, json, re, requests, tempfile, logging
from groq import Groq
from reportlab.lib.pagesizes import A4
from reportlab.platypus import *
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from datetime import datetime
from dotenv import load_dotenv
import pytesseract
from PIL import Image
import io

load_dotenv()

app = Flask(__name__)
client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# Configure logging
logging.basicConfig(level=logging.INFO)

# =========================
# OCR FUNCTION FOR IMAGES
# =========================
def extract_text_from_image(image_data):
    try:
        image = Image.open(io.BytesIO(image_data))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        logging.error(f"OCR error: {e}")
        return ""

# =========================
# ENHANCED EXTRACT WITH OCR
# =========================
@app.route('/extract', methods=['POST'])
def extract():
    file_path = request.json.get("filePath")
    slide_types_detected = set()

    try:
        if not file_path:
            return jsonify({"error": "No file path provided"}), 400

        if file_path.startswith("http"):
            response = requests.get(file_path)
            temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            temp.write(response.content)
            temp.close()
            prs = Presentation(temp.name)
        else:
            prs = Presentation(file_path)

        slides = []
        slide_keywords = {
            "Problem Statement": ["problem", "pain point", "challenge", "issue"],
            "Solution/Product": ["solution", "product", "platform", "app", "software"],
            "Market Opportunity": ["market", "tam", "sam", "som", "opportunity", "size"],
            "Business Model": ["revenue", "business model", "pricing", "subscription", "freemium"],
            "Competitive Landscape": ["competitor", "competition", "landscape", "vs", "advantage"],
            "Team": ["team", "founder", "ceo", "cto", "advisor", "leadership"],
            "Financial Projections": ["financial", "projection", "forecast", "revenue", "growth"],
            "Traction/Milestones": ["traction", "milestone", "users", "customers", "revenue", "growth"],
            "Funding Ask": ["funding", "investment", "raise", "ask", "use of funds"]
        }

        for slide_idx, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
                
                # OCR for images
                if hasattr(shape, "image") and shape.image:
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = extract_text_from_image(image_bytes)
                        text += ocr_text + " "
                    except:
                        pass
            
            if text.strip():
                slides.append(text.strip())
                
                # Detect slide types
                text_lower = text.lower()
                for slide_type, keywords in slide_keywords.items():
                    if any(keyword in text_lower for keyword in keywords):
                        slide_types_detected.add(slide_type)

        # Validate slide count
        if len(slides) < 5:
            return jsonify({"error": "Minimum 5 slides required"}), 400
        
        if len(slides) > 20:
            return jsonify({"error": "Maximum 20 slides allowed"}), 400
        
        # Validate required slide types (at least 3)
        if len(slide_types_detected) < 3:
            return jsonify({
                "error": f"Pitch deck must include at least 3 of these slide types: Problem Statement, Solution/Product, Market Opportunity, Business Model, Competitive Landscape, Team, Financial Projections, Traction/Milestones, Funding Ask. Detected: {list(slide_types_detected)}"
            }), 400

        return jsonify({
            "slides": slides,
            "detected_slide_types": list(slide_types_detected)
        })

    except Exception as e:
        logging.error(f"Extraction error: {e}")
        return jsonify({"error": str(e)}), 500

# =========================
# ENHANCED ANALYZE WITH WEIGHTS
# =========================
@app.route('/analyze', methods=['POST'])
def analyze():
    slides = request.json.get("slides", [])
    text = " ".join(slides)

    weights = {
        "Problem Statement": 10,
        "Solution/Product": 15,
        "Market Opportunity": 20,
        "Business Model": 15,
        "Competitive Landscape": 10,
        "Team": 15,
        "Traction/Milestones": 10,
        "Financial Projections": 10,
        "Clarity and Presentation": 5
    }

    prompt = f"""
Return ONLY valid JSON. No markdown, no explanations.

{{
 "startup_name": "",
 "recommendation": "Strong Buy | Hold | Pass",
 "confidence_score": 0,
 "weights": {json.dumps(weights)},
 "categories": {{
  "Problem Statement": {{"score": 0, "feedback": ""}},
  "Solution/Product": {{"score": 0, "feedback": ""}},
  "Market Opportunity": {{"score": 0, "feedback": ""}},
  "Business Model": {{"score": 0, "feedback": ""}},
  "Competitive Landscape": {{"score": 0, "feedback": ""}},
  "Team": {{"score": 0, "feedback": ""}},
  "Traction/Milestones": {{"score": 0, "feedback": ""}},
  "Financial Projections": {{"score": 0, "feedback": ""}},
  "Clarity and Presentation": {{"score": 0, "feedback": ""}}
 }},
 "strengths": [],
 "weaknesses": [],
 "recommendations": ""
}}

RULES:
- recommendation MUST be exactly: "Strong Buy", "Hold", or "Pass"
- Extract startup_name from slides (if not found, use "Unknown Startup")
- confidence_score: 0-100 based on data completeness (more data = higher score)
- Each category score: 0-10
- Each category feedback: 50-150 words
- strengths: 3-5 bullet points
- weaknesses: 3-5 bullet points
- recommendations: 100-200 words

Scoring Guidelines by Category:
1. Problem Statement (10% weight): Score based on clarity, evidence of customer pain, scope
2. Solution/Product (15% weight): Score based on feasibility, innovation, alignment
3. Market Opportunity (20% weight): Score based on TAM/SAM/SOM clarity, realism, data
4. Business Model (15% weight): Score based on revenue streams, scalability, pricing
5. Competitive Landscape (10% weight): Score based on competitor analysis, UVP strength
6. Team (15% weight): Score based on experience, completeness, track record
7. Traction/Milestones (10% weight): Score based on metrics, progress, achievements
8. Financial Projections (10% weight): Score based on forecasts, assumptions, realism
9. Clarity and Presentation (5% weight): Score based on flow, design, professionalism

Slides content:
{text[:15000]}
"""

    try:
        response = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama-3.1-8b-instant",
            temperature=0.3
        )

        raw = response.choices[0].message.content
        raw = re.sub(r"```json|```", "", raw).strip()

        try:
            result = json.loads(raw)
        except json.JSONDecodeError as e:
            logging.error(f"JSON parse error: {e}\nRaw: {raw}")
            return jsonify({"error": "AI returned invalid JSON"}), 500

        # Calculate overall score using weights
        total_weighted_score = 0
        for category, weight in weights.items():
            if category in result["categories"]:
                total_weighted_score += result["categories"][category]["score"] * weight
        
        result["overall_score"] = round(total_weighted_score / 10)  # Convert to 0-100 scale

        # Ensure confidence_score is set
        if "confidence_score" not in result or result["confidence_score"] == 0:
            # Calculate confidence based on how many categories have meaningful scores
            non_zero_categories = sum(1 for cat in result["categories"].values() if cat["score"] > 5)
            result["confidence_score"] = min(100, int((non_zero_categories / 9) * 100))

        return jsonify(result)

    except Exception as e:
        logging.error(f"Analysis error: {e}")
        return jsonify({"error": str(e)}), 500

# =========================
# ENHANCED PDF GENERATION
# =========================
@app.route('/generate-pdf', methods=['POST'])
def generate_pdf():
    data = request.json

    try:
        timestamp = datetime.utcnow().strftime("%d-%m-%Y %H:%M:%S UTC")
        startup_name = data.get("startup_name", "Startup").replace(" ", "_")
        date_str = datetime.now().strftime("%d%m%Y")
        file_name = f"Investment_Thesis_{startup_name}_{date_str}.pdf"

        # Create document with custom styles
        doc = SimpleDocTemplate(file_name, pagesize=A4, 
                                leftMargin=0.75*inch, rightMargin=0.75*inch,
                                topMargin=0.75*inch, bottomMargin=0.75*inch)
        
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=16,
            textColor='#1a472a',
            alignment=TA_CENTER,
            spaceAfter=30
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            textColor='#2d5a3f',
            spaceBefore=15,
            spaceAfter=10
        )
        
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontSize=11,
            leading=14,
            spaceAfter=6
        )

        elements = []

        # Title
        elements.append(Paragraph("Investment Thesis Report", title_style))
        elements.append(Spacer(1, 0.2*inch))

        # Summary Section
        elements.append(Paragraph("Executive Summary", heading_style))
        elements.append(Paragraph(f"<b>Startup:</b> {data.get('startup_name')}", normal_style))
        elements.append(Paragraph(f"<b>Recommendation:</b> {data.get('recommendation')}", normal_style))
        elements.append(Paragraph(f"<b>Overall Score:</b> {data.get('overall_score')}/100", normal_style))
        elements.append(Paragraph(f"<b>Confidence Score:</b> {data.get('confidence_score')}%", normal_style))
        elements.append(Paragraph(f"<b>Processing Date:</b> {timestamp}", normal_style))
        elements.append(Spacer(1, 0.2*inch))

        # Category Analysis
        elements.append(Paragraph("Detailed Category Analysis", heading_style))
        
        weights = data.get("weights", {
            "Problem Statement": 10, "Solution/Product": 15, "Market Opportunity": 20,
            "Business Model": 15, "Competitive Landscape": 10, "Team": 15,
            "Traction/Milestones": 10, "Financial Projections": 10, "Clarity and Presentation": 5
        })
        
        for category, values in data.get("categories", {}).items():
            weight = weights.get(category, 0)
            elements.append(Paragraph(
                f"<b>{category}</b> (Weight: {weight}%) - Score: {values['score']}/10",
                heading_style
            ))
            elements.append(Paragraph(values.get("feedback", "No feedback provided"), normal_style))
            elements.append(Spacer(1, 0.1*inch))

        # Strengths and Weaknesses
        cols = [['Strengths', 'Weaknesses']]
        data_table = [[
            Paragraph("<b>Strengths</b>", normal_style),
            Paragraph("<b>Weaknesses</b>", normal_style)
        ]]
        
        strengths_text = "<br/>".join([f"- {s}" for s in data.get("strengths", [])])
        weaknesses_text = "<br/>".join([f"- {w}" for w in data.get("weaknesses", [])])
        
        data_table.append([
            Paragraph(strengths_text if strengths_text else "None identified", normal_style),
            Paragraph(weaknesses_text if weaknesses_text else "None identified", normal_style)
        ])
        
        from reportlab.platypus import Table
        table = Table(data_table, colWidths=[2.5*inch, 2.5*inch])
        elements.append(table)
        elements.append(Spacer(1, 0.2*inch))

        # Recommendations
        elements.append(Paragraph("Recommendations", heading_style))
        elements.append(Paragraph(data.get("recommendations", "No recommendations provided"), normal_style))

        # Footer
        elements.append(Spacer(1, 0.5*inch))
        elements.append(Paragraph(f"Report generated by PitchDeck AI - {timestamp}", 
                                 ParagraphStyle('Footer', parent=normal_style, fontSize=8, alignment=TA_CENTER)))

        doc.build(elements)
        return send_file(file_name, as_attachment=True, download_name=file_name)

    except Exception as e:
        logging.error(f"PDF generation error: {str(e)}")
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(port=5001, debug=True, host='0.0.0.0')